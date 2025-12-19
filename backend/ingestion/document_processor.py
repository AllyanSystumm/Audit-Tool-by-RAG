"""Multi-format document processor for extracting text from various file types."""

import os
from pathlib import Path
from typing import Dict, Optional
import logging
import re
import unicodedata

# Document processing libraries
from docx import Document as DocxDocument
import docx2python
import fitz  # PyMuPDF
from pptx import Presentation
import openpyxl
import pandas as pd
from lxml import etree

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


class DocumentProcessor:
    """Extract text from multiple document formats."""
    
    SUPPORTED_FORMATS = {
        'docx': 'Word Document',
        'pdf': 'PDF Document',
        'pptx': 'PowerPoint Presentation',
        'xlsx': 'Excel Spreadsheet',
        'xls': 'Excel Spreadsheet (Legacy)',
        'xml': 'XML Document',
        'txt': 'Text Document'
    }
    
    def __init__(self):
        """Initialize the document processor."""
        self.extractors = {
            'docx': self._extract_from_docx,
            'pdf': self._extract_from_pdf,
            'pptx': self._extract_from_pptx,
            'xlsx': self._extract_from_xlsx,
            'xls': self._extract_from_xls,
            'xml': self._extract_from_xml,
            'txt': self._extract_from_txt
        }
    
    def is_supported(self, file_path: str) -> bool:
        """Check if file format is supported."""
        ext = Path(file_path).suffix.lower().lstrip('.')
        return ext in self.SUPPORTED_FORMATS
    
    def extract_text(self, file_path: str) -> Dict[str, any]:
        """
        Extract text from document.
        
        Args:
            file_path: Path to the document
            
        Returns:
            Dictionary with 'text', 'metadata', and 'status'
        """
        if not os.path.exists(file_path):
            return {
                'text': '',
                'metadata': {},
                'status': 'error',
                'error': 'File not found'
            }
        
        ext = Path(file_path).suffix.lower().lstrip('.')
        
        if ext not in self.extractors:
            return {
                'text': '',
                'metadata': {},
                'status': 'error',
                'error': f'Unsupported file format: .{ext}'
            }
        
        try:
            extractor = self.extractors[ext]
            text = extractor(file_path)
            text = self._clean_text(text)
            
            return {
                'text': text.strip(),
                'metadata': {
                    'filename': Path(file_path).name,
                    'format': ext,
                    'size': os.path.getsize(file_path),
                    'char_count': len(text)
                },
                'status': 'success',
                'error': None
            }
        except Exception as e:
            logger.error(f"Error extracting text from {file_path}: {str(e)}")
            return {
                'text': '',
                'metadata': {},
                'status': 'error',
                'error': str(e)
            }

    def _clean_text(self, text: str) -> str:
        """
        Normalize and sanitize extracted text to improve chunking/retrieval.
        Removes common Confluence export artifacts and HTML tags.
        """
        if not text:
            return ""

        # Unicode normalization (fix ligatures like "ﬁ" → "fi" that can appear as missing letters)
        text = unicodedata.normalize("NFKC", text)

        # Replace a few remaining common ligatures that sometimes survive extraction
        text = (
            text.replace("\ufb00", "ff")
            .replace("\ufb01", "fi")
            .replace("\ufb02", "fl")
            .replace("\ufb03", "ffi")
            .replace("\ufb04", "ffl")
            .replace("\ufb05", "ft")
            .replace("\ufb06", "st")
        )

        # Remove common Confluence/Atlassian export artifacts
        # e.g. "Error rendering macro 'toc': null"
        text = re.sub(r"(?im)^Error rendering macro 'toc'\\s*:\\s*null\\s*$", "", text)
        text = re.sub(r"(?i)Error rendering macro 'toc'\\s*:\\s*null", "", text)
        # Some exports show 'oc' (cropped/buggy macro name)
        text = re.sub(r"(?im)^Error rendering macro 'oc'\\s*:\\s*null\\s*$", "", text)
        text = re.sub(r"(?i)Error rendering macro 'oc'\\s*:\\s*null", "", text)

        # Strip HTML tags that sometimes appear in extracted DOCX content
        text = re.sub(r"<[^>]+>", " ", text)

        # Collapse excessive horizontal whitespace (but preserve newlines)
        # NOTE: Avoid \t/\f/\v here: in some regex engines/configs this can accidentally match literal letters.
        text = re.sub(r"[^\S\r\n]+", " ", text)
        text = re.sub(r"\\n{3,}", "\\n\\n", text)

        return text.strip()
    
    def _extract_from_docx(self, file_path: str) -> str:
        """Extract text from DOCX file."""
        # Prefer python-docx as primary extraction.
        # In some DOCX exports, docx2python can drop characters (e.g., "Verification" → "Veri ica ion").
        text_parts = []

        try:
            doc = DocxDocument(file_path)

            # Paragraphs
            for para in doc.paragraphs:
                if para.text and para.text.strip():
                    text_parts.append(para.text)

            # Tables
            for table in doc.tables:
                for row in table.rows:
                    row_cells = [cell.text.strip() for cell in row.cells if cell.text and cell.text.strip()]
                    if row_cells:
                        text_parts.append(" | ".join(row_cells))

        except Exception:
            # Fallback: docx2python (best-effort)
            try:
                docx_content = docx2python.docx2python(file_path)
                for page in docx_content.body:
                    for section in page:
                        for paragraph in section:
                            if isinstance(paragraph, list):
                                text_parts.append(' '.join(str(item) for item in paragraph if item))
                            else:
                                text_parts.append(str(paragraph))
            except Exception:
                pass

        return '\n'.join(text_parts)
    
    def _extract_from_pdf(self, file_path: str) -> str:
        """Extract text from PDF file using PyMuPDF."""
        text_parts = []
        
        with fitz.open(file_path) as doc:
            for page_num, page in enumerate(doc):
                text = page.get_text()
                if text.strip():
                    text_parts.append(f"--- Page {page_num + 1} ---\n{text}")
        
        return '\n\n'.join(text_parts)
    
    def _extract_from_pptx(self, file_path: str) -> str:
        """Extract text from PowerPoint file."""
        text_parts = []
        prs = Presentation(file_path)
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = [f"--- Slide {slide_num} ---"]
            
            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
                
                # Extract from tables
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text.strip():
                                row_text.append(cell.text)
                        if row_text:
                            slide_text.append(' | '.join(row_text))
            
            if len(slide_text) > 1:  # More than just the header
                text_parts.append('\n'.join(slide_text))
        
        return '\n\n'.join(text_parts)
    
    def _extract_from_xlsx(self, file_path: str) -> str:
        """Extract text from Excel XLSX file."""
        text_parts = []
        
        # Try pandas first (easier for structured data)
        try:
            xl_file = pd.ExcelFile(file_path)
            for sheet_name in xl_file.sheet_names:
                df = pd.read_excel(xl_file, sheet_name=sheet_name)
                text_parts.append(f"--- Sheet: {sheet_name} ---")
                text_parts.append(df.to_string(index=False))
        except:
            # Fallback to openpyxl
            wb = openpyxl.load_workbook(file_path, data_only=True)
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                text_parts.append(f"--- Sheet: {sheet_name} ---")
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = [str(cell) for cell in row if cell is not None]
                    if row_text:
                        text_parts.append(' | '.join(row_text))
        
        return '\n\n'.join(text_parts)
    
    def _extract_from_xls(self, file_path: str) -> str:
        """Extract text from Excel XLS file (legacy format)."""
        text_parts = []
        
        try:
            xl_file = pd.ExcelFile(file_path, engine='xlrd')
            for sheet_name in xl_file.sheet_names:
                df = pd.read_excel(xl_file, sheet_name=sheet_name)
                text_parts.append(f"--- Sheet: {sheet_name} ---")
                text_parts.append(df.to_string(index=False))
        except Exception as e:
            logger.error(f"Error reading XLS file: {e}")
            raise
        
        return '\n\n'.join(text_parts)
    
    def _extract_from_xml(self, file_path: str) -> str:
        """Extract text from XML file."""
        try:
            tree = etree.parse(file_path)
            root = tree.getroot()
            
            # Extract all text content
            text_parts = []
            for element in root.iter():
                if element.text and element.text.strip():
                    text_parts.append(element.text.strip())
                if element.tail and element.tail.strip():
                    text_parts.append(element.tail.strip())
            
            return '\n'.join(text_parts)
        except Exception as e:
            logger.error(f"Error parsing XML: {e}")
            # Fallback: read as plain text
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()
    
    def _extract_from_txt(self, file_path: str) -> str:
        """Extract text from plain text file."""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()


# Example usage
if __name__ == "__main__":
    processor = DocumentProcessor()
    
    # Test with a sample file
    result = processor.extract_text("sample.docx")
    print(f"Status: {result['status']}")
    print(f"Text length: {len(result['text'])}")
    print(f"Metadata: {result['metadata']}")

